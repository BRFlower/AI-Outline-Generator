"""
    该模块实现了函数extract_text_from_file，其接受一个文件路径，通过大模型提取文件的主要内容，输出与主要内容有关的文本
    extract_text_from_file:
        input:
            filepath: str #文件路径
            *api_key: str #dify项目的api key
            *user: str #用户名
        output: [str] #文件文本内容，流式输出，列表所有字符串相加即为最终结果
"""

import json
import requests
import os
import chardet
from pathlib import Path

API_KEY = "app-bnJeTo2iBB1ErxkA7kCBmPBv"

def _upload_file(file_path, api_key=API_KEY, user="noob"):
    url = 'https://api.dify.ai/v1/files/upload'
    headers = {
        'Authorization': f'Bearer {api_key}'
    }

    with open(file_path, 'rb') as file:
        files = {'file': (os.path.basename(file_path), file, 'text/plain')}
        payload = {'user': user}
        response = requests.post(url, headers=headers, files=files, data=payload)

    if response.status_code == 201:
        print("文件上传成功")
        return response.json()['id']
    else:
        print(f"文件上传失败，状态码：{response.status_code}，信息：{response.text}")
        return None


def _run_extraction(file_id, api_key=API_KEY, user="noob"):
    url = 'https://api.dify.ai/v1/workflows/run'
    headers = {'Authorization': f'Bearer {api_key}'}
    payload = {
        'inputs': {
            'file': {
                'transfer_method': 'local_file',
                'upload_file_id': file_id,
                'type': 'document'
            }
        },
        'user': user,
        'response_mode': 'streaming'
    }
    response = requests.post(url, headers=headers, json=payload)
    if response.status_code == 200:
        print("文本提取成功")
        stream = []
        for line in response.iter_lines():
            if line and line[:6] == b'data: ':
                line = line[6:]
                data = json.loads(line)
                content = data.get('data', {}).get('text', None)
                if content:
                    stream.append(content)
        return stream
    else:
        print(f"文本提取失败，状态码：{response.status_code}，信息：{response.text}")
        return [response.text]


def extract_text_from_file(filepath, api_key=API_KEY, user="noob"):
    file_id = _upload_file(filepath, api_key, user)
    return _run_extraction(file_id, api_key, user)

# ==================== 新增：传统文本提取功能 ====================

def detect_encoding(file_path):
    """检测文件编码"""
    with open(file_path, 'rb') as file:
        raw_data = file.read()
        result = chardet.detect(raw_data)
        encoding = result['encoding']
        confidence = result['confidence']
        
        # 如果检测置信度较低，尝试常见编码
        if confidence < 0.7:
            common_encodings = ['utf-8', 'gbk', 'gb2312', 'utf-8-sig', 'latin-1', 'cp1252']
            for enc in common_encodings:
                try:
                    raw_data.decode(enc)
                    encoding = enc
                    break
                except UnicodeDecodeError:
                    continue
        
        return encoding

def extract_text_from_text_file(file_path):
    """提取文本文件内容（支持各种文本格式）"""
    try:
        encoding = detect_encoding(file_path)
        with open(file_path, 'r', encoding=encoding, errors='ignore') as file:
            content = file.read()
        return [content]
    except Exception as e:
        print(f"读取文本文件失败: {str(e)}")
        return [f"读取文件失败: {str(e)}"]

def extract_text_from_docx(file_path):
    """提取docx文件文本"""
    try:
        from docx import Document
        doc = Document(file_path)
        content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content.append(paragraph.text)
        return ['\n'.join(content)]
    except ImportError:
        return ["请安装python-docx库: pip install python-docx"]
    except Exception as e:
        print(f"读取docx文件失败: {str(e)}")
        return [f"读取文件失败: {str(e)}"]

def extract_text_from_pptx(file_path):
    """提取pptx文件文本"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        content = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                content.append('\n'.join(slide_text))
        return ['\n\n'.join(content)]
    except ImportError:
        return ["请安装python-pptx库: pip install python-pptx"]
    except Exception as e:
        print(f"读取pptx文件失败: {str(e)}")
        return [f"读取文件失败: {str(e)}"]

def extract_text_from_pdf(file_path):
    """提取pdf文件文本"""
    try:
        import PyPDF2
        content = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text.strip():
                    content.append(text)
        return ['\n'.join(content)]
    except ImportError:
        return ["请安装PyPDF2库: pip install PyPDF2"]
    except Exception as e:
        print(f"读取pdf文件失败: {str(e)}")
        return [f"读取文件失败: {str(e)}"]

def extract_text_traditional(filepath):
    """
    使用传统方法提取文件文本内容
    支持格式：txt, md, docx, pptx, pdf, 以及cursor编辑器能打开的其他文本格式
    """
    if not os.path.exists(filepath):
        return [f"文件不存在: {filepath}"]
    
    file_path = Path(filepath)
    file_extension = file_path.suffix.lower()
    
    print(f"正在提取文件: {filepath}")
    print(f"文件类型: {file_extension}")
    
    # 支持的文本格式（cursor编辑器能直接打开的）
    text_extensions = {
        '.txt', '.md', '.markdown', '.rst', '.tex', '.html', '.htm', 
        '.xml', '.json', '.csv', '.tsv', '.log', '.ini', '.cfg', '.conf',
        '.py', '.js', '.ts', '.jsx', '.tsx', '.vue', '.java', '.cpp', '.c',
        '.h', '.hpp', '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.kt',
        '.scala', '.sql', '.sh', '.bat', '.ps1', '.yaml', '.yml', '.toml',
        '.properties', '.env', '.gitignore', '.dockerfile', '.makefile',
        '.cmake', '.sass', '.scss', '.less', '.css', '.styl', '.coffee',
        '.lua', '.pl', '.pm', '.r', '.m', '.f', '.f90', '.f95', '.f03',
        '.ada', '.pas', '.d', '.nim', '.zig', '.v', '.sv', '.vhd', '.vhdl'
    }
    
    # 根据文件扩展名选择提取方法
    if file_extension in text_extensions:
        return extract_text_from_text_file(filepath)
    elif file_extension == '.docx':
        return extract_text_from_docx(filepath)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(filepath)
    elif file_extension == '.pdf':
        return extract_text_from_pdf(filepath)
    else:
        # 对于未知格式，尝试作为文本文件读取
        print(f"未知文件格式 {file_extension}，尝试作为文本文件读取")
        return extract_text_from_text_file(filepath)

if __name__ == '__main__':
    # 测试代码
    test_files = [
        'example_material.txt',
        'example_material.docx',
        'example_material.pdf',
        'example_material.md',
        'example_material.pptx'
    ]
    
    for file_path in test_files:
        if os.path.exists(file_path):
            print(f"\n=== 测试文件: {file_path} ===")
            result = extract_text_traditional(file_path)
            print(f"提取结果长度: {len(''.join(result))} 字符")
            print(f"提取结果预览: {''.join(result)[:200]}...")
        else:
            print(f"测试文件不存在: {file_path}")